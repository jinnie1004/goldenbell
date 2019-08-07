from pptx import Presentation
import os
import sys

# question = [1,2,3,4,5]
# answer = [11,12,13,14,15]
# q_a_li = [[i, j] for i, j in zip(question, answer)]

def making_PPT(q_a_li):
    try:
        current_path = sys._MEIPASS
    except AttributeError:
        current_path = '.'
    prs = Presentation(os.path.join(current_path, '골든벨_템플릿.pptx'))
    title_slide_layout = prs.slide_layouts[0]
    for q_a in q_a_li:
        q, a = q_a
        slide = prs.slides.add_slide(title_slide_layout)

        for shape in slide.shapes:
            if shape.shape_id == 3:
                shape.text = f'{q}?'
            elif shape.shape_id == 2:
                shape.text = f"{a}"
            elif shape.shape_id == 4:
                shape.text = '문 제'
    prs.save('골든벨.pptx')

    os.startfile('골든벨.pptx')
